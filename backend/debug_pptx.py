"""
Debug script: parse a PPTX file and print raw structural info for theme colors,
slide background, layout shapes, master shapes, and slide shapes.

Usage:
    python debug_pptx.py [path/to/file.pptx]

If no path is given it looks for Business_2Col_Teal.pptx in Downloads,
then falls back to any .pptx it finds there.
"""
from __future__ import annotations

import glob
import os
import sys


# ── locate file ────────────────────────────────────────────────────────────────

def find_pptx() -> str | None:
    if len(sys.argv) > 1:
        path = sys.argv[1]
        if os.path.isfile(path):
            return path
        print(f"[WARN] Argument '{path}' not found, falling back to auto-search.")

    preferred = r"C:/Users/UserAdmin/Downloads/Business_2Col_Teal.pptx"
    if os.path.isfile(preferred):
        return preferred

    pattern = r"C:/Users/UserAdmin/Downloads/*.pptx"
    matches = glob.glob(pattern)
    if matches:
        return matches[0]

    return None


# ── theme color extraction ─────────────────────────────────────────────────────

def extract_theme_colors(prs) -> dict:
    """Extract the 12 standard theme color slots from the presentation theme XML."""
    theme_colors: dict[int, str] = {}
    try:
        from pptx.oxml.ns import qn
        from lxml import etree  # bundled with python-pptx

        # The theme element lives on the slide master
        for master in prs.slide_masters:
            try:
                theme_elem = master.element.find(
                    ".//" + qn("a:theme"), master.element.nsmap
                )
                if theme_elem is None:
                    # Try direct child path
                    theme_elem = master.theme_color_map  # may not exist
            except Exception:
                theme_elem = None

            # Walk the XML tree of the master part looking for clrScheme
            try:
                part = master.part
                # theme is a relationship target of the master part
                from pptx.opc.constants import RELATIONSHIP_TYPE as RT
                for rel in part.rels.values():
                    try:
                        target = rel.target_part
                        xml_bytes = target._blob
                        root = etree.fromstring(xml_bytes)
                        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                        clr_scheme = root.find(".//a:clrScheme", ns)
                        if clr_scheme is None:
                            continue
                        # Standard slot order:
                        SLOT_NAMES = [
                            "dk1", "lt1", "dk2", "lt2",
                            "accent1", "accent2", "accent3", "accent4",
                            "accent5", "accent6", "hlink", "folHlink",
                        ]
                        for idx, slot in enumerate(SLOT_NAMES, start=1):
                            elem = clr_scheme.find(f"a:{slot}", ns)
                            if elem is None:
                                continue
                            # Color value lives in child srgbClr or sysClr
                            rgb_elem = elem.find("a:srgbClr", ns)
                            if rgb_elem is not None:
                                val = rgb_elem.get("val", "")
                                if val:
                                    theme_colors[idx] = f"#{val.upper()}"
                                    continue
                            sys_elem = elem.find("a:sysClr", ns)
                            if sys_elem is not None:
                                last_clr = sys_elem.get("lastClr", "")
                                if last_clr:
                                    theme_colors[idx] = f"#{last_clr.upper()}"
                    except Exception:
                        pass
            except Exception:
                pass

            if theme_colors:
                break  # stop after first master that yields colors
    except Exception as exc:
        print(f"  [ERROR in extract_theme_colors] {exc}")

    return theme_colors


# ── shape info printer ─────────────────────────────────────────────────────────

def print_shape_info(shape, indent: str = "    ") -> None:
    name = "(unknown)"
    try:
        name = shape.name
    except Exception as e:
        name = f"<error: {e}>"

    shape_type = "(unknown)"
    try:
        shape_type = str(shape.shape_type)
    except Exception as e:
        shape_type = f"<error: {e}>"

    has_tf = "(unknown)"
    try:
        has_tf = str(shape.has_text_frame)
    except Exception as e:
        has_tf = f"<error: {e}>"

    fill_type = "(no fill attr)"
    try:
        fill_type = str(shape.fill.type)
    except Exception as e:
        fill_type = f"<error: {e}>"

    print(f"{indent}name={name!r}  shape_type={shape_type}  has_text_frame={has_tf}  fill.type={fill_type}")

    # If it has a text frame, print a snippet of text
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text[:120].replace("\n", "\\n")
            print(f"{indent}  text={text!r}")
    except Exception:
        pass


# ── main ───────────────────────────────────────────────────────────────────────

def main() -> None:
    pptx_path = find_pptx()
    if pptx_path is None:
        print("[ERROR] No .pptx file found. Pass a path as argument or place one in Downloads.")
        sys.exit(1)

    print(f"Parsing: {pptx_path}")
    print("=" * 70)

    try:
        from pptx import Presentation
    except ImportError:
        print("[ERROR] python-pptx is not installed. Run: pip install python-pptx")
        sys.exit(1)

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        print(f"[ERROR] Could not open presentation: {exc}")
        sys.exit(1)

    # ── 1. Theme colors ────────────────────────────────────────────────────────
    print("\n### 1. THEME COLORS ###")
    try:
        theme_colors = extract_theme_colors(prs)
        SLOT_LABELS = {
            1: "dk1 (dark1/text1)",
            2: "lt1 (light1/bg1)",
            3: "dk2 (dark2/text2)",
            4: "lt2 (light2/bg2)",
            5: "accent1",
            6: "accent2",
            7: "accent3",
            8: "accent4",
            9: "accent5",
            10: "accent6",
            11: "hlink",
            12: "folHlink",
        }
        if theme_colors:
            for slot, hex_val in sorted(theme_colors.items()):
                label = SLOT_LABELS.get(slot, f"slot{slot}")
                print(f"  [{slot:>2}] {label:<28} = {hex_val}")
        else:
            print("  (no theme colors extracted)")
    except Exception as exc:
        print(f"  [ERROR] {exc}")

    # ── 2. Slide 1 ─────────────────────────────────────────────────────────────
    if not prs.slides:
        print("\n[INFO] Presentation has no slides.")
        return

    slide = prs.slides[0]
    print(f"\n### 2. SLIDE 1 (of {len(prs.slides)}) ###")

    # 2a. Background fill type
    print("\n  -- slide.background.fill.type --")
    try:
        bg = slide.background
        fill = bg.fill
        print(f"  fill.type = {fill.type}")
        # Also try to read the fore color
        try:
            fc = fill.fore_color
            print(f"  fill.fore_color.type = {fc.type}")
            try:
                print(f"  fill.fore_color.rgb  = #{fc.rgb:06X}")
            except Exception as e:
                print(f"  fill.fore_color.rgb  = <error: {e}>")
            try:
                print(f"  fill.fore_color.theme_color = {fc.theme_color}")
            except Exception as e:
                print(f"  fill.fore_color.theme_color = <error: {e}>")
        except Exception as e:
            print(f"  fore_color = <error: {e}>")
    except Exception as exc:
        print(f"  [ERROR] {exc}")

    # 2b. Slide layout shapes
    print("\n  -- slide.slide_layout.shapes --")
    try:
        layout = slide.slide_layout
        shapes = layout.shapes
        print(f"  ({len(shapes)} shapes in layout '{layout.name}')")
        for i, shape in enumerate(shapes):
            print(f"  [{i}]", end="")
            print_shape_info(shape, indent=" ")
    except Exception as exc:
        print(f"  [ERROR] {exc}")

    # 2c. Slide master shapes
    print("\n  -- slide.slide_layout.slide_master.shapes --")
    try:
        master = slide.slide_layout.slide_master
        shapes = master.shapes
        print(f"  ({len(shapes)} shapes in master)")
        for i, shape in enumerate(shapes):
            print(f"  [{i}]", end="")
            print_shape_info(shape, indent=" ")
    except Exception as exc:
        print(f"  [ERROR] {exc}")

    # 2d. Slide shapes
    print("\n  -- slide.shapes --")
    try:
        shapes = slide.shapes
        print(f"  ({len(shapes)} shapes on slide)")
        for i, shape in enumerate(shapes):
            print(f"  [{i}]", end="")
            print_shape_info(shape, indent=" ")
    except Exception as exc:
        print(f"  [ERROR] {exc}")

    print("\n" + "=" * 70)
    print("Done.")


if __name__ == "__main__":
    main()
