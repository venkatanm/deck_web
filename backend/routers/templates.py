"""
Templates router — handles user and admin template management including PPTX import.

PPTX parsing uses python-pptx.  The conversion is best-effort:
  • Text boxes  → text elements
  • Rectangles / rounded-rects → rect elements
  • Pictures    → image placeholder elements (no actual image data)
  • Fill colors are extracted as hex strings

Install: pip install python-pptx  (added to requirements.txt)
"""
from __future__ import annotations

import io
import uuid
from typing import Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from pydantic import BaseModel
from sqlalchemy import select, or_
from sqlalchemy.ext.asyncio import AsyncSession

from auth import get_current_user
from database import get_db
from models import Template, User

router = APIRouter(prefix="/api/templates", tags=["templates"])

MAX_PPTX_SIZE = 50 * 1024 * 1024  # 50 MB


# ── helpers ────────────────────────────────────────────────────────────────────

def _color_to_hex(color, theme_colors: Optional[dict] = None) -> Optional[str]:
    """Convert pptx RGBColor / ThemeColor to hex string, returns None if unavailable."""
    if color is None:
        return None
    try:
        # Direct RGB color — RGBColor doesn't support format spec, use str()
        rgb_str = str(color.rgb).upper()
        if len(rgb_str) == 6:
            return f"#{rgb_str}"
    except Exception:
        pass
    try:
        from pptx.dml.color import _NoneColor
        if isinstance(color, _NoneColor):
            return None
    except Exception:
        pass
    try:
        # Try to resolve theme color via the presentation theme palette
        tc = int(color.theme_color)
        # Apply luminance modifier if present (tints/shades)
        # luminance_modulation is on the underlying _color object or readable from XML
        lum_mod = None
        lum_off = None
        try:
            from pptx.oxml.ns import qn as _qn
            # Access underlying _color object then its XML element
            _c = getattr(color, '_color', None) or color
            _xclr = getattr(_c, '_xClr', None)
            if _xclr is not None:
                lm_el = _xclr.find(_qn('a:lumMod'))
                lo_el = _xclr.find(_qn('a:lumOff'))
                if lm_el is not None:
                    lum_mod = int(lm_el.get('val', '100000')) / 100000.0
                if lo_el is not None:
                    lum_off = int(lo_el.get('val', '0')) / 100000.0
        except Exception:
            pass
        # Fallback: try direct attribute (works for _RGBColor in some pptx versions)
        if lum_mod is None:
            try:
                lum_mod = color.luminance_modulation / 100000.0
            except Exception:
                pass
        if lum_off is None:
            try:
                lum_off = color.luminance_offset / 100000.0
            except Exception:
                pass

        # Use extracted theme palette if available
        if theme_colors and tc in theme_colors:
            hex_c = theme_colors[tc]
        else:
            # Fallback palette (matches Office default theme)
            THEME_MAP = {
                1: "#FFFFFF",  # background 1
                2: "#000000",  # text 1
                3: "#FFFFFF",  # background 2
                4: "#000000",  # text 2
                5: "#4472C4",  # accent 1
                6: "#ED7D31",  # accent 2
                7: "#A9D18E",  # accent 3
                8: "#FF0000",  # accent 4
                9: "#FFC000",  # accent 5
                10: "#5B9BD5", # accent 6
            }
            hex_c = THEME_MAP.get(tc, "#333333")

        # Apply luminance modulation/offset (tint/shade)
        if (lum_mod is not None or lum_off is not None) and hex_c:
            try:
                r = int(hex_c[1:3], 16)
                g = int(hex_c[3:5], 16)
                b = int(hex_c[5:7], 16)
                # Convert to HLS, apply mod/off, convert back
                import colorsys
                h, l, s = colorsys.rgb_to_hls(r/255, g/255, b/255)
                if lum_mod is not None:
                    l = l * lum_mod
                if lum_off is not None:
                    l = l + lum_off
                l = max(0.0, min(1.0, l))
                r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
                hex_c = f"#{round(r2*255):02X}{round(g2*255):02X}{round(b2*255):02X}"
            except Exception:
                pass
        return hex_c
    except Exception:
        return None


def _extract_theme_colors(prs) -> dict:
    """
    Extract the theme color palette (dk1, lt1, dk2, lt2, accent1-6) from the
    presentation theme XML. Returns a dict mapping theme_color int → hex string.
    Theme color ints follow pptx MSO_THEME_COLOR enum:
      1=dk1/text1, 2=lt1/bg1, 3=dk2/text2, 4=lt2/bg2, 5..10=accent1..6
    """
    result = {}
    slot_order = ["dk1", "lt1", "dk2", "lt2",
                  "accent1", "accent2", "accent3",
                  "accent4", "accent5", "accent6"]

    def _parse_clr_scheme(clrScheme, qn):
        for idx, slot in enumerate(slot_order, start=1):
            el = clrScheme.find(qn(f"a:{slot}"))
            if el is None:
                continue
            rgb_el = el.find(qn("a:srgbClr"))
            if rgb_el is not None:
                val = rgb_el.get("val", "")
                if len(val) == 6:
                    result[idx] = f"#{val.upper()}"
            else:
                sys_el = el.find(qn("a:sysClr"))
                if sys_el is not None:
                    last = sys_el.get("lastClr", "")
                    if len(last) == 6:
                        result[idx] = f"#{last.upper()}"

    try:
        from pptx.oxml.ns import qn
        from lxml import etree

        # Strategy 1: find theme via slide master rels (theme is a separate Part)
        for rel in prs.slide_masters[0].part.rels.values():
            if "theme" in rel.reltype.lower():
                try:
                    blob = rel._target.blob
                    theme_xml = etree.fromstring(blob)
                    clrScheme = theme_xml.find(f".//{qn('a:clrScheme')}")
                    if clrScheme is not None:
                        _parse_clr_scheme(clrScheme, qn)
                        if result:
                            break  # found theme, stop searching rels
                except Exception:
                    pass

        # Strategy 2: look for a:clrScheme anywhere in master element tree
        try:
            master_el = prs.slide_masters[0].element
            clrScheme = master_el.find(f".//{qn('a:clrScheme')}")
            if clrScheme is not None:
                _parse_clr_scheme(clrScheme, qn)
        except Exception:
            pass
    except Exception:
        pass

    # Add MSO alias mappings: BACKGROUND_1=14→lt1(2), TEXT_1=13→dk1(1),
    # TEXT_2=15→dk2(3), BACKGROUND_2=16→lt2(4)
    ALIAS_MAP = {13: 1, 14: 2, 15: 3, 16: 4}
    for alias, canonical in ALIAS_MAP.items():
        if alias not in result and canonical in result:
            result[alias] = result[canonical]

    return result


def _resolve_fill_color(fill, default: Optional[str] = "#cccccc", theme_colors: Optional[dict] = None) -> Optional[str]:
    """Resolve a pptx FillElement to a hex color string. Returns None if fill is transparent/invisible."""
    try:
        from pptx.oxml.ns import qn
        if fill is None or fill.type is None:
            return default
        # Solid fill
        try:
            c = _color_to_hex(fill.fore_color, theme_colors)
            if c:
                return c
        except Exception:
            pass
        # Solid fill via raw XML schemeClr (handles bg1/bg2/tx1/tx2 aliases)
        try:
            solid_el = fill._xPr.find(qn('a:solidFill'))
            if solid_el is not None:
                scheme_el = solid_el.find(qn('a:schemeClr'))
                if scheme_el is not None:
                    SCHEME_ALIAS = {
                        "bg1": 14, "lt1": 2,
                        "bg2": 16, "lt2": 4,
                        "tx1": 13, "dk1": 1,
                        "tx2": 15, "dk2": 3,
                        "accent1": 5, "accent2": 6, "accent3": 7,
                        "accent4": 8, "accent5": 9, "accent6": 10,
                    }
                    val = scheme_el.get("val", "")
                    idx = SCHEME_ALIAS.get(val)
                    if idx and theme_colors and idx in theme_colors:
                        import colorsys as _cs
                        hex_c = theme_colors[idx]
                        lm_el = scheme_el.find(qn('a:lumMod'))
                        lo_el = scheme_el.find(qn('a:lumOff'))
                        if lm_el is not None or lo_el is not None:
                            lm = int(lm_el.get('val','100000'))/100000 if lm_el is not None else 1.0
                            lo = int(lo_el.get('val','0'))/100000 if lo_el is not None else 0.0
                            r,g,b = int(hex_c[1:3],16),int(hex_c[3:5],16),int(hex_c[5:7],16)
                            h,l,s = _cs.rgb_to_hls(r/255,g/255,b/255)
                            l = max(0.0, min(1.0, l*lm + lo))
                            r2,g2,b2 = _cs.hls_to_rgb(h,l,s)
                            hex_c = f"#{round(r2*255):02X}{round(g2*255):02X}{round(b2*255):02X}"
                        return hex_c
        except Exception:
            pass
        # Gradient — find first non-transparent stop; if ALL stops are transparent, return None
        try:
            stops = fill._xPr.findall(f".//{qn('a:gs')}")
            if stops:
                found_opaque = False
                for stop in stops:
                    clr_el = stop.find(qn('a:schemeClr')) or stop.find(qn('a:srgbClr'))
                    if clr_el is not None:
                        alpha_el = clr_el.find(qn('a:alpha'))
                        if alpha_el is not None and int(alpha_el.get('val', '100000')) == 0:
                            continue
                        from pptx.dml.color import ColorFormat
                        try:
                            cf = ColorFormat.from_colorelement(clr_el)
                            c = _color_to_hex(cf, theme_colors)
                            if c:
                                found_opaque = True
                                return c
                        except Exception:
                            pass
                if not found_opaque:
                    return None  # all stops transparent — skip shape
        except Exception:
            pass
    except Exception:
        pass
    return default


def _emu_to_px(emu: int, canvas_px: int, slide_emu: int) -> float:
    """Scale EMU coordinates to our canvas pixel coordinates."""
    if slide_emu == 0:
        return 0
    return round(emu / slide_emu * canvas_px, 2)


def _shape_fill_color(shape, default: Optional[str] = "#cccccc", theme_colors: Optional[dict] = None) -> Optional[str]:
    """Extract fill color from a shape. Returns None if fill is invisible/transparent."""
    try:
        return _resolve_fill_color(shape.fill, default, theme_colors)
    except Exception:
        return default


def _shape_stroke(shape, theme_colors: Optional[dict] = None):
    """Extract stroke color and width from a shape. Returns (color_or_None, width_px)."""
    try:
        line = shape.line
        if line and line.width and line.width > 0:
            width_pt = max(round(line.width.pt), 1)
            color = None
            try:
                color = _color_to_hex(line.color, theme_colors)
            except Exception:
                pass
            return color or "#000000", width_pt
    except Exception:
        pass
    return None, 0


def _slide_bg_color(slide, theme_colors: Optional[dict] = None) -> Optional[str]:
    """Extract per-slide background color, falling back to layout then master."""
    from pptx.enum.dml import MSO_THEME_COLOR

    def _fill_color(obj):
        try:
            from pptx.oxml.ns import qn
            SCHEME_MAP = {
                'dk1': 1, 'lt1': 2, 'dk2': 3, 'lt2': 4,
                'bg1': 14, 'bg2': 16, 'tx1': 13, 'tx2': 15,
                'accent1': 5, 'accent2': 6, 'accent3': 7,
                'accent4': 8, 'accent5': 9, 'accent6': 10,
            }
            ALIAS_MAP = {13: 1, 14: 2, 15: 3, 16: 4}

            def _resolve_scheme(val):
                idx = SCHEME_MAP.get(val)
                if idx and theme_colors:
                    resolved = ALIAS_MAP.get(idx, idx)
                    return theme_colors.get(resolved)
                return None

            def _luminance(h):
                try:
                    hx = h.lstrip('#')
                    r, g, b = int(hx[0:2],16), int(hx[2:4],16), int(hx[4:6],16)
                    return 0.2126*r + 0.7152*g + 0.0722*b
                except Exception:
                    return 255

            # Access bgPr XML directly (background fills use p:bg/p:bgPr, not p:spPr)
            bg_el = obj.background._element
            p_bg = bg_el.find(qn('p:bg')) if bg_el is not None else None
            bg_pr = p_bg.find(qn('p:bgPr')) if p_bg is not None else None
            if bg_pr is not None:
                # Solid fill
                solid = bg_pr.find(qn('a:solidFill'))
                if solid is not None:
                    scheme = solid.find(qn('a:schemeClr'))
                    srgb = solid.find(qn('a:srgbClr'))
                    if scheme is not None:
                        c = _resolve_scheme(scheme.get('val', ''))
                        if c:
                            return c
                    if srgb is not None:
                        val = srgb.get('val', '')
                        if len(val) == 6:
                            return f'#{val.upper()}'

                # Gradient fill — collect opaque stops
                opaque = []
                for stop in bg_pr.findall(f'.//{qn("a:gs")}'):
                    _sc = stop.find(qn('a:schemeClr'))
                    _sr = stop.find(qn('a:srgbClr'))
                    clr_el = _sc if _sc is not None else _sr
                    if clr_el is None:
                        continue
                    alpha_el = clr_el.find(qn('a:alpha'))
                    if alpha_el is not None and int(alpha_el.get('val','100000')) == 0:
                        continue
                    val = clr_el.get('val', '')
                    c = _resolve_scheme(val) or (f'#{val.upper()}' if len(val) == 6 else None)
                    if c:
                        opaque.append(c)
                if opaque:
                    return min(opaque, key=_luminance)  # pick darkest stop

            # Fallback: use pptx fill API
            fill = obj.background.fill
            if fill.type is not None:
                c = _resolve_fill_color(fill, "#cccccc", theme_colors)
                if c and c != "#cccccc":
                    return c
        except Exception:
            pass
        return None

    # 1. Slide-specific background
    c = _fill_color(slide)
    if c:
        return c
    # 2. Layout background
    try:
        c = _fill_color(slide.slide_layout)
        if c:
            return c
    except Exception:
        pass
    # 3. Master background
    try:
        c = _fill_color(slide.slide_layout.slide_master)
        if c:
            return c
    except Exception:
        pass
    return None


def _font_family_map(pptx_name: Optional[str]) -> str:
    """Map a PPTX font name to the closest available web font."""
    if not pptx_name:
        return "Inter"
    name_lower = pptx_name.lower().strip()
    # Theme font placeholders
    if name_lower in ("+mj-lt", "+mn-lt", "calibri light", "calibri"):
        return "Inter"
    mapping = {
        "arial": "Inter",
        "helvetica": "Inter",
        "trebuchet ms": "Inter",
        "verdana": "Inter",
        "tahoma": "Inter",
        "gill sans": "Inter",
        "gill sans mt": "Inter",
        "century gothic": "Montserrat",
        "futura": "Montserrat",
        "montserrat": "Montserrat",
        "poppins": "Poppins",
        "open sans": "Open Sans",
        "lato": "Lato",
        "raleway": "Raleway",
        "oswald": "Oswald",
        "times new roman": "Georgia",
        "times": "Georgia",
        "georgia": "Georgia",
        "garamond": "Playfair Display",
        "palatino": "Playfair Display",
        "playfair display": "Playfair Display",
        "merriweather": "Merriweather",
        "courier": "Inter",
        "courier new": "Inter",
        "impact": "Oswald",
        "bebas neue": "Bebas Neue",
    }
    for key, val in mapping.items():
        if key in name_lower:
            return val
    # Return original name — browser will fall back to Inter if unavailable
    return pptx_name


def _bg_is_dark(bg_hex: str) -> bool:
    """Return True if the hex background color is perceptually dark."""
    try:
        h = bg_hex.lstrip("#")
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        # Relative luminance (sRGB approximation)
        lum = 0.2126 * r + 0.7152 * g + 0.0722 * b
        return lum < 128
    except Exception:
        return False


def _parse_text_frame(tf, base: dict, slide_w_emu: int, slide_h_emu: int, canvas_w: int, canvas_h: int, theme_colors: Optional[dict] = None, bg_color: str = "#ffffff") -> Optional[dict]:
    """Parse a text frame into a text element. Returns None if empty."""
    from pptx.enum.text import PP_ALIGN

    paragraphs_text = []
    font_size = 18
    fill_color = "#ffffff" if _bg_is_dark(bg_color) else "#1e293b"
    align = "left"
    bold = False
    italic = False
    font_family = "Inter"

    for para in tf.paragraphs:
        para_text = "".join(run.text for run in para.runs)
        paragraphs_text.append(para_text)

        if para.runs:
            run = para.runs[0]
            if run.font.size:
                try:
                    font_size = round(run.font.size.pt)
                except Exception:
                    pass
            if run.font.color and run.font.color.type is not None:
                c = _color_to_hex(run.font.color, theme_colors)
                if c:
                    fill_color = c
            if run.font.bold:
                bold = True
            if run.font.italic:
                italic = True
            try:
                fn = run.font.name
                if fn:
                    font_family = _font_family_map(fn)
            except Exception:
                pass
        try:
            if para.alignment == PP_ALIGN.CENTER:
                align = "center"
            elif para.alignment == PP_ALIGN.RIGHT:
                align = "right"
        except Exception:
            pass

    full_text = "\n".join(paragraphs_text).strip()
    if not full_text:
        return None

    font_style = "normal"
    if bold and italic:
        font_style = "bold italic"
    elif bold:
        font_style = "bold"
    elif italic:
        font_style = "italic"

    return {
        **base,
        "type": "text",
        "content": full_text,
        "text": full_text,
        "fontSize": max(font_size, 8),
        "fontFamily": font_family,
        "fontStyle": font_style,
        "fill": fill_color,
        "align": align,
    }


def _get_placeholder_font_size(ph_shape, ph_type: str) -> int:
    """
    Walk the placeholder inheritance chain to find font size:
    slide shape → layout placeholder (same idx) → master placeholder → master txStyles.
    Returns pt size as int (default 28 for title, 18 for body).
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    def _sz_from_xml_el(el) -> Optional[int]:
        """Find first a:rPr or a:defPPr with sz attribute in a txBody XML element."""
        if el is None:
            return None
        for tag in (qn("a:defPPr"), qn("a:lvl1pPr"), qn("a:rPr")):
            for rpr in el.iter(tag):
                sz = rpr.get("sz")
                if sz:
                    try:
                        return round(int(sz) / 100)
                    except Exception:
                        pass
        return None

    # 1. Try the slide shape itself
    try:
        for para in ph_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return round(run.font.size.pt)
    except Exception:
        pass

    # 2. Try layout placeholder with same idx
    try:
        idx = ph_shape.placeholder_format.idx
        layout = ph_shape.part.slide_layout
        for layout_ph in layout.placeholders:
            if layout_ph.placeholder_format.idx == idx:
                for para in layout_ph.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            return round(run.font.size.pt)
                # Walk layout shape txBody XML
                txBody = layout_ph._element.find(qn("p:txBody"))
                sz = _sz_from_xml_el(txBody)
                if sz:
                    return sz
    except Exception:
        pass

    # 3. Try master txStyles for title / body
    try:
        master = ph_shape.part.slide_layout.slide_master
        txStyles = master.element.find(qn("p:txStyles"))
        if txStyles is not None:
            if ph_type == "title":
                title_style = txStyles.find(qn("p:titleStyle"))
                sz = _sz_from_xml_el(title_style)
                if sz:
                    return sz
            else:
                body_style = txStyles.find(qn("p:bodyStyle"))
                sz = _sz_from_xml_el(body_style)
                if sz:
                    return sz
    except Exception:
        pass

    # 4. Defaults
    return 36 if ph_type == "title" else 20


def _get_placeholder_font_color(ph_shape, bg_hex: str, theme_colors: dict) -> str:
    """
    Get text color for a placeholder. Check run → layout ph → theme dk1/dk2,
    with automatic white fallback on dark backgrounds.
    """
    dark_bg = _bg_is_dark(bg_hex)
    default_color = "#ffffff" if dark_bg else "#1e293b"

    # 1. Direct color from runs on the slide shape
    try:
        for para in ph_shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.color and run.font.color.type is not None:
                    c = _color_to_hex(run.font.color, theme_colors)
                    if c:
                        return c
    except Exception:
        pass

    # 2. Theme dk1 color for light bg, or white for dark bg
    if dark_bg:
        return "#ffffff"
    # Use dk1 (index 1) from theme
    return theme_colors.get(1, "#1e293b")


def _render_slide_thumbnail(slide, canvas_w: int, canvas_h: int) -> Optional[str]:
    """
    Render a slide as a PNG thumbnail and return as base64 data URL.
    Uses LibreOffice or python-pptx thumbnail if available, otherwise returns None.
    """
    try:
        # Try using the slide's thumbnail if python-pptx supports it
        # (pptx >= 0.6.21 has Slide.shapes thumbnail support via libreoffice)
        pass
    except Exception:
        pass
    return None


def _render_layout_as_image(slide, canvas_w: int, canvas_h: int, slide_w_emu: int, slide_h_emu: int, theme_colors: dict, bg_color: str) -> Optional[str]:
    """
    Render the layout/master decorative shapes (non-placeholder) as a PNG image.
    Returns base64 data URL or None if PIL unavailable.
    """
    try:
        from PIL import Image as PILImage, ImageDraw
        import base64
        from pptx.oxml.ns import qn

        def _hex_to_rgba(hex_c: Optional[str], alpha: int = 255):
            if not hex_c:
                return None
            h = hex_c.lstrip("#")
            if len(h) == 6:
                return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16), alpha)
            return None

        # Start with the slide background color as the base
        bg_rgba = _hex_to_rgba(bg_color, 255) or (255, 255, 255, 255)
        img = PILImage.new("RGBA", (canvas_w, canvas_h), bg_rgba)
        draw = ImageDraw.Draw(img)

        def _draw_shape(shape):
            try:
                left   = round(_emu_to_px(shape.left   or 0, canvas_w, slide_w_emu))
                top    = round(_emu_to_px(shape.top    or 0, canvas_h, slide_h_emu))
                width  = round(_emu_to_px(shape.width  or 0, canvas_w, slide_w_emu))
                height = round(_emu_to_px(shape.height or 0, canvas_h, slide_h_emu))
                if width <= 0 or height <= 0:
                    return

                # Picture shape — paste image onto canvas
                try:
                    from pptx.enum.shapes import MSO_SHAPE_TYPE
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        pic_img = PILImage.open(io.BytesIO(shape.image.blob)).convert("RGBA")
                        pic_img = pic_img.resize((width, height), PILImage.LANCZOS)
                        img.paste(pic_img, (left, top), pic_img)
                        return
                except Exception:
                    pass

                fill_c = _resolve_fill_color(shape.fill, None, theme_colors)
                if fill_c is None:
                    return

                rgba = _hex_to_rgba(fill_c)
                if rgba is None:
                    return

                draw.rectangle([left, top, left + width, top + height], fill=rgba)
            except Exception:
                pass

        # Draw master shapes first (deepest layer), then layout shapes on top
        try:
            for shape in slide.slide_layout.slide_master.shapes:
                try:
                    if shape.is_placeholder:
                        continue
                except Exception:
                    pass
                _draw_shape(shape)
        except Exception:
            pass

        try:
            for shape in slide.slide_layout.shapes:
                try:
                    if shape.is_placeholder:
                        continue
                except Exception:
                    pass
                _draw_shape(shape)
        except Exception:
            pass

        if img.getbbox() is None:
            return None  # blank image — no decoration found

        # Convert to PNG data URL
        out_buf = io.BytesIO()
        img.save(out_buf, format="PNG", optimize=True)
        b64 = base64.b64encode(out_buf.getvalue()).decode("utf-8")
        return f"data:image/png;base64,{b64}"
    except Exception:
        return None


def _parse_pptx(data: bytes) -> tuple[list[dict], int, int, str]:
    """
    Parse a PPTX file using hybrid approach:
    - Render master/layout decorative shapes as a background image element
    - Parse placeholder shapes (title, body, picture) as editable canvas elements
    - Font sizes derived by walking slide → layout → master inheritance chain
    - Produces faithful visual representation with fully editable text

    Returns (slides, canvas_w, canvas_h, first_bg_color)
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    except ImportError:
        raise HTTPException(500, "python-pptx is not installed on the server.")

    prs = Presentation(io.BytesIO(data))
    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    theme_colors = _extract_theme_colors(prs)

    # Canvas size from aspect ratio
    aspect = slide_w_emu / slide_h_emu if slide_h_emu else 16 / 9
    canvas_h = 720
    canvas_w = round(canvas_h * aspect)

    # Padding for text elements (px)
    TEXT_PAD = 12

    slides_out = []
    first_bg = "#ffffff"

    for slide_idx, slide in enumerate(prs.slides):
        elements = []

        # Background color: slide → layout → master → white
        bg_color = _slide_bg_color(slide, theme_colors) or "#ffffff"
        if slide_idx == 0:
            first_bg = bg_color

        # Render layout/master decorative shapes as a background image
        layout_img_src = _render_layout_as_image(
            slide, canvas_w, canvas_h, slide_w_emu, slide_h_emu, theme_colors, bg_color
        )
        if layout_img_src:
            elements.append({
                "id": str(uuid.uuid4()),
                "type": "image",
                "src": layout_img_src,
                "imageId": None,
                "x": 0, "y": 0,
                "width": canvas_w, "height": canvas_h,
                "rotation": 0, "opacity": 1,
                "objectFit": "fill",
                "locked": True,  # decorative — user shouldn't accidentally move it
            })

        # Collect only placeholder shapes from the slide
        for shape in slide.shapes:
            try:
                if not shape.is_placeholder:
                    continue
            except Exception:
                continue

            try:
                ph_fmt = shape.placeholder_format
                ph_idx = ph_fmt.idx
                ph_type_raw = ph_fmt.type

                left   = _emu_to_px(shape.left   or 0, canvas_w, slide_w_emu)
                top    = _emu_to_px(shape.top    or 0, canvas_h, slide_h_emu)
                width  = _emu_to_px(shape.width  or 0, canvas_w, slide_w_emu)
                height = _emu_to_px(shape.height or 0, canvas_h, slide_h_emu)

                base = {
                    "id": str(uuid.uuid4()),
                    "x": left,
                    "y": top,
                    "width": max(width, 20),
                    "height": max(height, 20),
                    "rotation": getattr(shape, "rotation", 0) or 0,
                    "opacity": 1,
                }

                # Classify placeholder
                try:
                    from pptx.enum.shapes import PP_PLACEHOLDER as PP_PH
                    is_title = ph_idx == 0 or ph_type_raw in (
                        PP_PH.TITLE, PP_PH.CENTER_TITLE, PP_PH.VERTICAL_TITLE,
                    )
                    is_body = ph_idx in (1, 2) or ph_type_raw in (
                        PP_PH.BODY, PP_PH.OBJECT, PP_PH.SUBTITLE,
                    )
                    is_picture = ph_type_raw in (PP_PH.PICTURE,)
                except Exception:
                    is_title = ph_idx == 0
                    is_body = ph_idx in (1, 2)
                    is_picture = False

                # Picture placeholder → styled rect
                if is_picture or shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            import base64
                            img_blob = shape.image.blob
                            ct = shape.image.content_type or "image/png"
                            try:
                                from PIL import Image as PILImage
                                MAX_PX = 800
                                pil_img = PILImage.open(io.BytesIO(img_blob))
                                if pil_img.width > MAX_PX or pil_img.height > MAX_PX:
                                    pil_img.thumbnail((MAX_PX, MAX_PX), PILImage.LANCZOS)
                                out_buf = io.BytesIO()
                                fmt = "JPEG" if ct in ("image/jpeg", "image/jpg") else "PNG"
                                if fmt == "JPEG" and pil_img.mode in ("RGBA", "P"):
                                    pil_img = pil_img.convert("RGB")
                                pil_img.save(out_buf, format=fmt, quality=82, optimize=True)
                                img_blob = out_buf.getvalue()
                                ct = "image/jpeg" if fmt == "JPEG" else "image/png"
                            except Exception:
                                pass
                            b64 = base64.b64encode(img_blob).decode("utf-8")
                            elements.append({
                                **base,
                                "type": "image",
                                "src": f"data:{ct};base64,{b64}",
                                "imageId": None,
                                "objectFit": "cover",
                            })
                        else:
                            # Empty picture placeholder → dashed image drop zone
                            elements.append({
                                **base,
                                "type": "rect",
                                "fill": "#e2e8f0",
                                "stroke": "#94a3b8",
                                "strokeWidth": 2,
                                "cornerRadius": 4,
                                "dash": [6, 4],
                            })
                    except Exception:
                        pass
                    continue

                # Text placeholder (title or body)
                if not shape.has_text_frame:
                    continue

                # Get text content from slide shape
                tf = shape.text_frame
                paragraphs_text = []
                for para in tf.paragraphs:
                    para_text = "".join(run.text for run in para.runs)
                    paragraphs_text.append(para_text)
                full_text = "\n".join(paragraphs_text).strip()
                if not full_text:
                    continue

                ph_type_str = "title" if is_title else "body"
                font_size = _get_placeholder_font_size(shape, ph_type_str)
                font_color = _get_placeholder_font_color(shape, bg_color, theme_colors)

                bold = is_title  # titles bold by default
                italic = False
                font_family = "Inter"
                align = "left"

                # Try to read explicit font properties from slide shape
                try:
                    from pptx.enum.text import PP_ALIGN
                    for para in tf.paragraphs:
                        if para.runs:
                            run = para.runs[0]
                            if run.font.bold is not None:
                                bold = run.font.bold
                            if run.font.italic:
                                italic = True
                            try:
                                fn = run.font.name
                                if fn:
                                    font_family = _font_family_map(fn)
                            except Exception:
                                pass
                        try:
                            if para.alignment == PP_ALIGN.CENTER:
                                align = "center"
                            elif para.alignment == PP_ALIGN.RIGHT:
                                align = "right"
                        except Exception:
                            pass
                        break  # only check first para for style
                except Exception:
                    pass

                font_style = "normal"
                if bold and italic:
                    font_style = "bold italic"
                elif bold:
                    font_style = "bold"
                elif italic:
                    font_style = "italic"

                # Scale font size to canvas
                # PPTX font sizes are in pt; at 720px canvas height (vs ~540pt slide),
                # scale proportionally so text fits visually
                scale = canvas_h / (slide_h_emu / 914400 * 72)  # slide height in pt
                scaled_font_size = max(round(font_size * scale), 8)

                elements.append({
                    **base,
                    "type": "text",
                    "content": full_text,
                    "text": full_text,
                    "fontSize": scaled_font_size,
                    "fontFamily": font_family,
                    "fontStyle": font_style,
                    "fill": font_color,
                    "align": align,
                    "padding": TEXT_PAD,
                })

            except Exception:
                pass  # skip broken placeholders

        slides_out.append({"elements": elements, "backgroundColor": bg_color})

    return _json_safe(slides_out), canvas_w, canvas_h, first_bg


def _json_safe(obj):
    """Recursively convert all values to JSON-serializable Python primitives."""
    if isinstance(obj, dict):
        return {k: _json_safe(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple)):
        return [_json_safe(v) for v in obj]
    if isinstance(obj, bool):
        return obj
    if isinstance(obj, int):
        return int(obj)
    if isinstance(obj, float):
        return float(obj)
    if obj is None:
        return None
    if isinstance(obj, str):
        return obj
    # Fallback: convert unknown types (numpy floats, Decimal, etc.) via str then float/int
    try:
        return float(obj)
    except Exception:
        pass
    return str(obj)


def _bg_color(prs) -> str:
    """Try to extract background color from the first slide; fallback to white."""
    try:
        slide = prs.slides[0]
        c = _slide_bg_color(slide)
        if c:
            return c
    except Exception:
        pass
    return "#ffffff"


# ── endpoints ─────────────────────────────────────────────────────────────────

@router.get("")
async def list_templates(
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """Return global templates + this user's private templates."""
    result = await db.execute(
        select(Template).where(
            or_(Template.is_global.is_(True), Template.user_id == user.id)
        ).order_by(Template.created_at.desc())
    )
    templates = result.scalars().all()
    return [_template_out(t) for t in templates]


@router.post("/upload-pptx")
async def upload_pptx(
    file: UploadFile = File(...),
    name: str = Form(""),
    is_global: bool = Form(False),
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    """
    Upload a PPTX file, parse it, and store as a Template.
    is_global=True requires admin privileges.
    """
    if is_global and not user.is_admin:
        raise HTTPException(403, "Admin privileges required to create global templates.")

    data = await file.read()
    if len(data) > MAX_PPTX_SIZE:
        raise HTTPException(413, "File exceeds 50 MB limit.")
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(400, "Only .pptx files are supported.")

    # Parse PPTX
    try:
        slides, canvas_w, canvas_h, bg = _parse_pptx(data)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(422, f"Could not parse PPTX: {e}")

    template_name = name.strip() or file.filename.replace(".pptx", "").replace("_", " ").title()

    # Upsert: if a same-named template already exists in the same scope, replace it
    owner_id = None if is_global else user.id
    existing = await db.execute(
        select(Template).where(
            Template.name == template_name,
            Template.user_id == owner_id,
            Template.is_global == is_global,
        )
    )
    tpl = existing.scalar_one_or_none()
    if tpl:
        tpl.canvas_size = {"width": canvas_w, "height": canvas_h, "backgroundColor": bg}
        tpl.slides = slides
        tpl.thumbnail_bg = bg
        tpl.source_format = "pptx"
    else:
        tpl = Template(
            user_id=owner_id,
            name=template_name,
            category="Global" if is_global else "My Uploads",
            tags=["pptx", "uploaded"],
            canvas_size={"width": canvas_w, "height": canvas_h, "backgroundColor": bg},
            slides=slides,
            thumbnail_bg=bg,
            is_global=is_global,
            source_format="pptx",
        )
        db.add(tpl)
    await db.commit()
    await db.refresh(tpl)
    return _template_out(tpl)


@router.delete("/{template_id}")
async def delete_template(
    template_id: str,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    result = await db.execute(
        select(Template).where(Template.id == template_id)
    )
    tpl = result.scalar_one_or_none()
    if not tpl:
        raise HTTPException(404, "Template not found.")
    # Non-admins can only delete their own templates
    if not user.is_admin and tpl.user_id != user.id:
        raise HTTPException(403, "Not authorized.")
    await db.delete(tpl)
    await db.commit()
    return {"ok": True}


# ── Admin: promote user to admin ──────────────────────────────────────────────

class PromoteBody(BaseModel):
    email: str

@router.post("/admin/promote")
async def promote_to_admin(
    body: PromoteBody,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    if not user.is_admin:
        raise HTTPException(403, "Admin privileges required.")
    result = await db.execute(select(User).where(User.email == body.email))
    target = result.scalar_one_or_none()
    if not target:
        raise HTTPException(404, "User not found.")
    target.is_admin = True
    await db.commit()
    return {"ok": True, "email": target.email}


@router.get("/admin/users")
async def list_users(
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
):
    if not user.is_admin:
        raise HTTPException(403, "Admin privileges required.")
    result = await db.execute(select(User).order_by(User.created_at.desc()))
    users = result.scalars().all()
    return [
        {"id": str(u.id), "email": u.email, "is_admin": u.is_admin,
         "created_at": u.created_at.isoformat()}
        for u in users
    ]


# ── helper ────────────────────────────────────────────────────────────────────

def _template_out(t: Template) -> dict:
    return {
        "id": str(t.id),
        "name": t.name,
        "category": t.category,
        "tags": t.tags or [],
        "canvasSize": t.canvas_size,
        "canvas_size": t.canvas_size,
        "slides": t.slides or [],
        "thumbnail": {"bg": t.thumbnail_bg or "#f1f5f9"},
        "thumbnail_bg": t.thumbnail_bg,
        "is_global": t.is_global,
        "source_format": t.source_format,
        "created_at": t.created_at.isoformat() if t.created_at else None,
    }
